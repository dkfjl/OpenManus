from __future__ import annotations

import asyncio
from pathlib import Path
from typing import List, Optional

from app.config import config
from app.exceptions import ToolError
from app.tool.base import BaseTool


class PptxPresentationTool(BaseTool):
    """Tool for creating/updating a PowerPoint .pptx presentation within the workspace.

    Accepts either `sections` (heading/content/bullets like the Word tool) or
    `slides` (title/body/bullets). A title slide is created automatically when
    `presentation_title` is provided.
    """

    name: str = "pptx_presentation"
    description: str = (
        "Create or update a .pptx presentation using structured sections or slides."
    )
    parameters: dict = {
        "type": "object",
        "properties": {
            "filepath": {
                "type": "string",
                "description": "Target .pptx path (relative to workspace or absolute within it).",
            },
            "presentation_title": {
                "type": "string",
                "description": "Optional title used for the first title slide.",
            },
            "sections": {
                "type": "array",
                "description": "Sections to convert into slides (heading/content/bullets).",
                "items": {
                    "type": "object",
                    "properties": {
                        "heading": {"type": "string"},
                        "content": {"type": "string"},
                        "bullets": {
                            "type": "array",
                            "items": {"type": "string"},
                        },
                    },
                    "required": [],
                },
            },
            "slides": {
                "type": "array",
                "description": "Explicit slide definitions (title/body/bullets)",
                "items": {
                    "type": "object",
                    "properties": {
                        "title": {"type": "string"},
                        "body": {"type": "string"},
                        "bullets": {
                            "type": "array",
                            "items": {"type": "string"},
                        },
                        "images": {
                            "type": "array",
                            "description": "Optional image paths to add to the slide",
                            "items": {"type": "string"},
                        },
                        "table": {
                            "type": "object",
                            "description": "Optional table to render on this slide",
                            "properties": {
                                "headers": {
                                    "type": "array",
                                    "items": {"type": "string"},
                                    "description": "Header row (optional)",
                                },
                                "rows": {
                                    "type": "array",
                                    "items": {
                                        "type": "array",
                                        "items": {"type": "string"},
                                    },
                                    "description": "Table body rows",
                                },
                            },
                            "required": ["rows"],
                        },
                    },
                    "required": [],
                },
            },
            "default_image_width": {
                "type": "number",
                "description": "Default image width in inches when adding pictures.",
                "default": 6.0,
            },
            "append": {
                "type": "boolean",
                "description": "If true and the file exists, append slides instead of recreating.",
                "default": False,
            },
            "author": {
                "type": "string",
                "description": "Optional author metadata stored in presentation properties.",
            },
        },
        "required": ["filepath"],
    }

    async def execute(
        self,
        *,
        filepath: str,
        presentation_title: Optional[str] = None,
        sections: Optional[List[dict]] = None,
        slides: Optional[List[dict]] = None,
        append: bool = False,
        author: Optional[str] = None,
        default_image_width: float = 6.0,
        **_: str,
    ):
        if not sections and not slides:
            raise ToolError("Provide `sections` or `slides` to create the presentation.")

        target_path = self._resolve_path(filepath)
        target_path.parent.mkdir(parents=True, exist_ok=True)

        # Defer import to keep dependency optional unless PPTX is requested
        try:
            from pptx import Presentation  # type: ignore
            from pptx.util import Pt, Inches
        except Exception as e:  # pragma: no cover - environment dependent
            raise ToolError(
                f"python-pptx is required for PPTX export ({e}). Please install 'python-pptx'."
            )

        def _write() -> dict:
            prs = Presentation(str(target_path)) if (append and target_path.exists()) else Presentation()

            slides_written = 0

            if presentation_title and not append:
                # Title slide (layout 0)
                try:
                    title_slide_layout = prs.slide_layouts[0]
                    slide = prs.slides.add_slide(title_slide_layout)
                    if slide.shapes.title:
                        slide.shapes.title.text = presentation_title
                    # Try to set subtitle if placeholder exists
                    if len(slide.placeholders) > 1:
                        subtitle = slide.placeholders[1]
                        subtitle.text = ""
                    slides_written += 1
                except Exception:
                    # Fallback to title-and-content layout if 0 is unavailable
                    layout = prs.slide_layouts[1]
                    slide = prs.slides.add_slide(layout)
                    slide.shapes.title.text = presentation_title
                    slides_written += 1

            def add_content_slide(title: str, body: Optional[str], bullets_list: Optional[List[str]]):
                layout = prs.slide_layouts[1]  # Title and Content
                slide = prs.slides.add_slide(layout)
                if slide.shapes.title:
                    slide.shapes.title.text = (title or "").strip() or "(无标题)"
                # Content placeholder (usually index 1)
                text_frame = None
                for shape in slide.placeholders:
                    if getattr(shape, "has_text_frame", False) and shape.placeholder_format.idx != 0:
                        text_frame = shape.text_frame
                        break
                if text_frame is None:
                    # Create a text frame in the first available shape with text_frame
                    for shape in slide.shapes:
                        if getattr(shape, "has_text_frame", False):
                            text_frame = shape.text_frame
                            break
                if text_frame is None:
                    return
                # Clear and write content
                text_frame.clear()

                def add_paragraph(line: str, level: int = 0):
                    if len(text_frame.paragraphs) == 0:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()
                    p.level = max(0, min(5, level))
                    run = p.add_run()
                    run.text = line
                    try:
                        run.font.size = Pt(18)
                    except Exception:
                        pass

                # Write body text (split by blank lines into bullets)
                if body:
                    blocks = [b.strip() for b in body.replace("\r\n", "\n").split("\n\n") if b.strip()]
                    for block in blocks:
                        for line in block.split("\n"):
                            if line.strip():
                                add_paragraph(line.strip(), level=0)

                # Write bullets
                for b in (bullets_list or []):
                    text = (b or "").strip()
                    if text:
                        add_paragraph(text, level=0)

                return slide

            def add_images(slide, image_paths: List[str]):
                # Place images in simple vertical flow with default width
                from pptx.enum.shapes import MSO_SHAPE_TYPE
                left_margin = Inches(0.8)
                top_margin = Inches(1.8)
                cur_top = top_margin
                for p in image_paths:
                    img_path = _resolve_inside_workspace(p)
                    if not img_path.exists():
                        continue
                    try:
                        pic = slide.shapes.add_picture(str(img_path), left_margin, cur_top, width=Inches(max(1.0, float(default_image_width))))
                        cur_top = pic.top + pic.height + Inches(0.2)
                    except Exception:
                        # Skip invalid image
                        continue

            def add_table_slide(title: str, headers: Optional[List[str]], rows: List[List[str]]):
                layout = prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[1]  # Title Only preferred
                slide = prs.slides.add_slide(layout)
                if slide.shapes.title:
                    slide.shapes.title.text = (title or "表格").strip()
                slide_width = prs.slide_width
                slide_height = prs.slide_height
                left = Inches(0.6)
                top = Inches(1.8)
                width = slide_width - Inches(1.2)
                height = slide_height - Inches(2.2)

                nrows = len(rows) + (1 if headers else 0)
                ncols = max(len(headers or []), max((len(r) for r in rows), default=1))
                table_shape = slide.shapes.add_table(nrows, ncols, left, top, width, height)
                table = table_shape.table
                # Write headers
                r_idx = 0
                if headers:
                    for c, text in enumerate(headers):
                        table.cell(0, c).text = str(text)
                    r_idx = 1
                # Write body
                for row in rows:
                    for c, text in enumerate(row):
                        table.cell(r_idx, c).text = str(text)
                    r_idx += 1

                return slide

            def _resolve_inside_workspace(p: str) -> Path:
                base = config.workspace_root
                candidate = Path(p).expanduser()
                if not candidate.is_absolute():
                    candidate = base / candidate
                return candidate.resolve()

            # Prefer `slides`, otherwise map `sections` -> slides
            if slides:
                for s in slides:
                    slide = add_content_slide(
                        title=(s.get("title") or ""),
                        body=s.get("body"),
                        bullets_list=s.get("bullets"),
                    )
                    # images
                    imgs = s.get("images") or []
                    if imgs:
                        add_images(slide, imgs)
                    # table
                    table_def = s.get("table")
                    if table_def and isinstance(table_def, dict):
                        # For table-only slide, create a new slide with table for clarity
                        add_table_slide(
                            title=(s.get("title") or "表格"),
                            headers=table_def.get("headers"),
                            rows=table_def.get("rows") or [],
                        )
                    slides_written += 1
            elif sections:
                for sec in sections:
                    slide = add_content_slide(
                        title=(sec.get("heading") or ""),
                        body=sec.get("content"),
                        bullets_list=sec.get("bullets"),
                    )
                    # Optional images/tables in sections as well
                    imgs = sec.get("images") or []
                    if imgs:
                        add_images(slide, imgs)
                    table_def = sec.get("table")
                    if table_def and isinstance(table_def, dict):
                        add_table_slide(
                            title=(sec.get("heading") or "表格"),
                            headers=table_def.get("headers"),
                            rows=table_def.get("rows") or [],
                        )
                    slides_written += 1

            if author:
                try:
                    prs.core_properties.author = author
                except Exception:
                    pass

            prs.save(str(target_path))
            return {"slides_written": slides_written}

        result = await asyncio.to_thread(_write)
        result["path"] = str(target_path)
        result["mode"] = "append" if (append and Path(target_path).exists()) else "overwrite"
        return self.success_response(result)

    @staticmethod
    def _resolve_path(filepath: str) -> Path:
        base = config.workspace_root.resolve()
        candidate = Path(filepath).expanduser()
        if not candidate.is_absolute():
            candidate = base / candidate
        resolved = candidate.resolve()
        if resolved.suffix.lower() != ".pptx":
            raise ToolError("Only .pptx files are supported by pptx_presentation tool.")
        if base not in resolved.parents and resolved != base:
            raise ToolError(f"Target path {resolved} is outside of the workspace directory.")
        return resolved
