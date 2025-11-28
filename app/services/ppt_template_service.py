from __future__ import annotations

import asyncio
import re
import tempfile
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from app.config import config
from app.logger import logger


def _sanitize_filename(topic: str) -> str:
    sanitized = re.sub(r"[^\w\u4e00-\u9fff]+", "_", topic).strip("_") or "presentation"
    return f"{sanitized}.pptx"


def _resolve_workspace_path(target: Optional[str], fallback_name: str) -> Path:
    base = config.workspace_root
    candidate = Path(target) if target else (Path("reports") / fallback_name)
    if not candidate.is_absolute():
        candidate = base / candidate
    candidate = candidate.with_suffix(".pptx")
    candidate.parent.mkdir(parents=True, exist_ok=True)
    return candidate


async def inspect_ppt_template(template_path: Path) -> Dict[str, Any]:
    """Return layouts and placeholders information from a PPTX template.

    Response shape:
    {
      "layouts": [
         {"index": 0, "name": "Title Slide", "placeholders": [
            {"idx": 0, "type": "TITLE", "name": "Title 1"}, ...
         ]}, ...
      ]
    }
    """

    def _inspect_sync() -> Dict[str, Any]:
        from pptx import Presentation  # type: ignore

        prs = Presentation(str(template_path))
        layouts: List[Dict[str, Any]] = []
        for i, layout in enumerate(prs.slide_layouts):
            placeholders = []
            try:
                for ph in layout.placeholders:
                    try:
                        ptype = getattr(getattr(ph, "placeholder_format", None), "type", None)
                        pidx = getattr(getattr(ph, "placeholder_format", None), "idx", None)
                        placeholders.append(
                            {
                                "idx": int(pidx) if pidx is not None else None,
                                "type": getattr(ptype, "name", str(ptype)),
                                "name": getattr(ph, "name", ""),
                            }
                        )
                    except Exception:
                        continue
            except Exception:
                placeholders = []
            layouts.append(
                {
                    "index": i,
                    "name": getattr(layout, "name", f"Layout {i}"),
                    "placeholders": placeholders,
                }
            )
        return {"layouts": layouts}

    return await asyncio.to_thread(_inspect_sync)


def _select_title_placeholder(placeholders) -> Optional[Any]:
    try:
        from pptx.enum.shapes import PP_PLACEHOLDER

        # Prefer TITLE/CENTER_TITLE
        for ph in placeholders:
            pf = getattr(ph, "placeholder_format", None)
            if not pf:
                continue
            if pf.type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                return ph
        # Fallback to first text placeholder
        for ph in placeholders:
            if getattr(ph, "has_text_frame", False):
                return ph
    except Exception:
        # Fallback without enums
        for ph in placeholders:
            if getattr(ph, "has_text_frame", False):
                return ph
    return None


def _select_subtitle_placeholder(placeholders) -> Optional[Any]:
    try:
        from pptx.enum.shapes import PP_PLACEHOLDER

        for ph in placeholders:
            pf = getattr(ph, "placeholder_format", None)
            if pf and pf.type == PP_PLACEHOLDER.SUBTITLE:
                return ph
        # Fallback to next text placeholder that's not the first (title)
        seen_title = False
        for ph in placeholders:
            if getattr(ph, "has_text_frame", False):
                if not seen_title:
                    seen_title = True
                    continue
                return ph
    except Exception:
        # Fallback without enums
        seen_title = False
        for ph in placeholders:
            if getattr(ph, "has_text_frame", False):
                if not seen_title:
                    seen_title = True
                    continue
                return ph
    return None


def _select_body_placeholder(placeholders) -> Optional[Any]:
    try:
        from pptx.enum.shapes import PP_PLACEHOLDER

        for ph in placeholders:
            pf = getattr(ph, "placeholder_format", None)
            if pf and pf.type in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.CONTENT):
                return ph
        # Fallback to any text-frame placeholder except the likely title
        seen_title = False
        for ph in placeholders:
            if getattr(ph, "has_text_frame", False):
                if not seen_title:
                    seen_title = True
                    continue
                return ph
    except Exception:
        seen_title = False
        for ph in placeholders:
            if getattr(ph, "has_text_frame", False):
                if not seen_title:
                    seen_title = True
                    continue
                return ph
    return None


async def generate_ppt_from_template(
    *,
    template_path: Path,
    data: Dict[str, Any],
    output_path: Optional[str] = None,
) -> Dict[str, Any]:
    """Generate a presentation using a template and structured slide data.

    data schema:
    {"slides": [{
        "type": "cover"|"content",
        "layout_index": 0,
        "title": "...",
        "subtitle": "...",
        "content": ["bullet 1", "bullet 2"] | "plain text"
    }, ...]}
    """

    def _write_sync() -> Tuple[str, int]:
        from pptx import Presentation  # type: ignore

        slides_def: List[dict] = list(data.get("slides") or [])
        # Fall back filename from first cover or first slide title
        topic = None
        for s in slides_def:
            if s.get("title"):
                topic = s["title"]
                if s.get("type", "").lower() == "cover":
                    break
        abs_out = _resolve_workspace_path(output_path, _sanitize_filename(topic or "presentation")).resolve()

        prs = Presentation(str(template_path))
        written = 0

        for s in slides_def:
            stype = (s.get("type") or "content").lower()
            idx = s.get("layout_index")
            # Sensible defaults: 0 for cover, 1 for content
            if idx is None:
                idx = 0 if stype == "cover" else 1
            try:
                layout = prs.slide_layouts[int(idx)]
            except Exception:
                # Fallback: use 0/1
                layout = prs.slide_layouts[0 if stype == "cover" else 1]

            slide = prs.slides.add_slide(layout)
            placeholders = list(slide.placeholders)

            if stype == "cover":
                title_ph = _select_title_placeholder(placeholders)
                subtitle_ph = _select_subtitle_placeholder(placeholders)
                if title_ph and getattr(title_ph, "has_text_frame", False):
                    title_ph.text_frame.clear()
                    # Keep style by only writing text
                    p = title_ph.text_frame.paragraphs[0]
                    p.text = str(s.get("title") or "")
                if subtitle_ph and getattr(subtitle_ph, "has_text_frame", False):
                    subtitle_ph.text_frame.clear()
                    p = subtitle_ph.text_frame.paragraphs[0]
                    p.text = str(s.get("subtitle") or "")
                written += 1
                continue

            # content slide
            title_ph = _select_title_placeholder(placeholders)
            body_ph = _select_body_placeholder(placeholders)

            if title_ph and getattr(title_ph, "has_text_frame", False):
                title_ph.text_frame.clear()
                p = title_ph.text_frame.paragraphs[0]
                p.text = str(s.get("title") or "")

            # Write body
            if body_ph and getattr(body_ph, "has_text_frame", False):
                tf = body_ph.text_frame
                tf.clear()

                def add_para(text: str, level: int = 0):
                    # Preserve default style; do not override fonts
                    if len(tf.paragraphs) == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.level = max(0, min(5, int(level)))
                    # Use runs to inherit style
                    run = getattr(p, "runs", None)
                    if run is not None and len(p.runs) > 0:
                        p.runs[0].text = text
                    else:
                        p.text = text

                content = s.get("content")
                if isinstance(content, list):
                    for item in content:
                        text = str(item)
                        if not text.strip():
                            continue
                        add_para(text, level=0)
                elif isinstance(content, str) and content.strip():
                    for line in content.replace("\r\n", "\n").split("\n"):
                        if line.strip():
                            add_para(line.strip(), level=0)
            written += 1

        prs.save(str(abs_out))
        return str(abs_out), written

    out, count = await asyncio.to_thread(_write_sync)
    logger.info(f"PPT generated at {out}, slides={count}")
    return {"filepath": out, "slides_written": count}

